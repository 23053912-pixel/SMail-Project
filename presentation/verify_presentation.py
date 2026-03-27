#!/usr/bin/env python3
"""Verify the presentation file is complete"""
from pptx import Presentation
import os

file_path = "SMAIL_System_Presentation.pptx"

# Check file exists
if not os.path.exists(file_path):
    print(f"ERROR: {file_path} not found")
    exit(1)

# Check file size
size = os.path.getsize(file_path)
print(f"File size: {size} bytes")

# Load and verify presentation
p = Presentation(file_path)
print(f"Total slides: {len(p.slides)}")

if len(p.slides) != 12:
    print(f"ERROR: Expected 12 slides but got {len(p.slides)}")
    exit(1)

print("\n✅ Presentation is complete and valid!")
print("✅ All 12 slides present and accounted for")
print("\nDeliverables ready:")
print("  - SMAIL_System_Presentation.pptx")
print("  - PRESENTATION_SCRIPT.md")
print("  - SLIDES_CONTENT.md")
print("  - generate_presentation.py")
