#!/usr/bin/env python3
"""
SMAIL System - Automatic PowerPoint Generator
Generates a 12-slide presentation automatically
Run: python3 generate_presentation.py
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

def add_title_slide(prs, title, subtitle):
    """Add a title slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(26, 115, 232)  # Google Blue
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1.5))
    title_frame = title_box.text_frame
    title_frame.text = title
    title_frame.paragraphs[0].font.size = Pt(60)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
    title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Subtitle
    subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(9), Inches(1))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = subtitle
    subtitle_frame.paragraphs[0].font.size = Pt(28)
    subtitle_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
    subtitle_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

def add_content_slide(prs, title, content_list):
    """Add a content slide with title and bullet points"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(255, 255, 255)
    
    # Add title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.text = title
    title_frame.paragraphs[0].font.size = Pt(44)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = RGBColor(26, 115, 232)
    
    # Add content
    content_box = slide.shapes.add_textbox(Inches(0.7), Inches(1.3), Inches(8.6), Inches(5.5))
    text_frame = content_box.text_frame
    text_frame.word_wrap = True
    
    for i, line in enumerate(content_list):
        if i == 0:
            p = text_frame.paragraphs[0]
        else:
            p = text_frame.add_paragraph()
        
        p.text = line
        p.level = 0
        p.font.size = Pt(18)
        p.font.color.rgb = RGBColor(32, 33, 36)
        p.space_before = Pt(6)
        p.space_after = Pt(6)

def main():
    """Generate the presentation"""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    print("🎬 Creating SMAIL System Presentation...")
    
    # Slide 1: Title
    print("📌 Slide 1: Title")
    add_title_slide(prs, "SMAIL SYSTEM", 
                    "Email Categorization & Spam Detection Platform")
    
    # Slide 2: Problem Statement
    print("📌 Slide 2: Problem Statement")
    add_content_slide(prs, "The Challenge: Email Overload", [
        "🔴 Problems:",
        "• 100+ emails received daily",
        "• Manual categorization is time-consuming",
        "• Unreliable spam detection",
        "• No intelligent organization system",
        "",
        "✅ Solution: SMAIL SYSTEM",
        "• Automated email categorization",
        "• AI-powered spam detection",
        "• Intelligent folder organization"
    ])
    
    # Slide 3: Architecture
    print("📌 Slide 3: System Architecture")
    add_content_slide(prs, "System Architecture", [
        "Frontend Layer",
        "• Web interface with HTML, CSS, JavaScript",
        "",
        "API Layer",
        "• Express.js REST API (Port 3000)",
        "",
        "Backend Services",
        "• Email Fetching • Categorization • Spam Detection",
        "",
        "Data Layer",
        "• SQLite with indexing & WAL mode",
        "",
        "ML Engine",
        "• Python-based spam detection (Port 5001)"
    ])
    
    # Slide 4: Features
    print("📌 Slide 4: Key Features")
    add_content_slide(prs, "Key Features & Capabilities", [
        "📧 Email Management",
        "• Fetch from Gmail API • Auto-categorize • Mark important",
        "",
        "🤖 Spam Detection",
        "• ML-powered classification • Real-time analysis • 95%+ accuracy",
        "",
        "👤 User Experience",
        "• Secure OAuth 2.0 authentication • Session management",
        "• Customizable settings • Responsive interface"
    ])
    
    # Slide 5: Tech Stack
    print("📌 Slide 5: Technology Stack")
    add_content_slide(prs, "Technology Stack", [
        "🎨 Frontend: HTML5, CSS3, JavaScript ES6+",
        "",
        "⚙️ Backend: Node.js & Express.js",
        "",
        "💾 Database: SQLite (lightweight, fast, indexed)",
        "",
        "🧠 AI/ML: Python, Scikit-learn, Batch Processing",
        "",
        "🔗 External: Gmail API v1, Google OAuth 2.0"
    ])
    
    # Slide 6: Backend Services
    print("📌 Slide 6: Modular Backend Services")
    add_content_slide(prs, "Four Specialized Backend Services", [
        "1️⃣ Email Fetching Service",
        "   Gmail API integration, message parsing, rate limiting",
        "",
        "2️⃣ Categorization Service",
        "   Rule-based classification, 4 categories, batch processing",
        "",
        "3️⃣ Spam Detection Service",
        "   ML integration, probability scoring, error handling",
        "",
        "4️⃣ Email Orchestrator",
        "   Coordinates all services, full pipeline, data persistence"
    ])
    
    # Slide 7: Frontend UI
    print("📌 Slide 7: Frontend Interface")
    add_content_slide(prs, "User Interface & Screens", [
        "🔐 Authentication",
        "   Google Sign-In via OAuth 2.0",
        "",
        "📬 Email List",
        "   Folder navigation, email preview, spam flags",
        "",
        "📖 Email Details",
        "   Full content, From/To/Subject, spam score",
        "",
        "✍️ Compose",
        "   New email, send via Gmail API, draft saving"
    ])
    
    # Slide 8: ML Spam Detection
    print("📌 Slide 8: AI-Powered Spam Detection")
    add_content_slide(prs, "Machine Learning Spam Detection", [
        "📊 Model Performance",
        "   • Accuracy: 95%+ • Real-time predictions",
        "",
        "🔄 Detection Process",
        "   1. Extract email text → 2. Send to ML API → 3. Receive probability",
        "   4. Calculate score → 5. Flag if suspicious",
        "",
        "🎯 Hybrid Approach",
        "   • Machine Learning (70%) • Rule-Based (30%)",
        "   • Combined for optimal accuracy"
    ])
    
    # Slide 9: Performance
    print("📌 Slide 9: Performance Optimization")
    add_content_slide(prs, "Speed & Efficiency Improvements", [
        "⚡ Performance Gains",
        "   Load 50 emails: 50-60s → 15-25s (2-4x faster)",
        "   Refresh (cached): 30-40s → 100ms (300x faster)",
        "   Response size: 500KB → 80KB (85% smaller)",
        "",
        "✅ Optimization Techniques",
        "   • Database indexes • Gzip compression • Client caching",
        "   • Batch processing • Connection pooling • Smart retries"
    ])
    
    # Slide 10: Database
    print("📌 Slide 10: Database Design")
    add_content_slide(prs, "Database Schema & Optimization", [
        "💾 Tables",
        "   • emails: id, sender, recipient, subject, body, date, labels",
        "   • spam_results: email_id, is_spam, spam_score, probability",
        "",
        "🔍 Indexes",
        "   • sender, date, labels, read • spam_score, is_spam",
        "",
        "⚡ Optimizations",
        "   • WAL mode for concurrent reads • 10MB cache",
        "   • Batch inserts • PRAGMA synchronous=NORMAL"
    ])
    
    # Slide 11: API
    print("📌 Slide 11: REST API Endpoints")
    add_content_slide(prs, "RESTful API Endpoints", [
        "📧 Email Operations",
        "   GET /api/user • POST /api/emails/process • GET /api/emails/:folder",
        "",
        "🎯 Specialized",
        "   POST /api/emails/categorize • POST /api/emails/spam-detect",
        "   POST /api/send • POST /api/draft",
        "",
        "🔐 Authentication",
        "   JWT tokens • Bearer scheme • OAuth 2.0 integration",
        "",
        "📦 Response Format",
        "   Consistent JSON with success flag, data, statistics"
    ])
    
    # Slide 12: Roadmap
    print("📌 Slide 12: Deployment & Future Roadmap")
    add_content_slide(prs, "Deployment & Future Vision", [
        "✅ Current Status",
        "   • Fully functional development system • Production-ready code",
        "",
        "🚀 Future Roadmap",
        "   • Q2 2026: Cloud deployment • Q3 2026: Mobile app",
        "   • Q4 2026: Analytics dashboard • 2027: AI assistant",
        "",
        "📊 Key Metrics",
        "   • 4 modular services • 95%+ spam accuracy",
        "   • 15-25s load time • Scalable microservices architecture"
    ])
    
    # Save presentation
    output_file = "SMAIL_System_Presentation.pptx"
    prs.save(output_file)
    print(f"\n✅ Presentation created: {output_file}")
    print(f"📊 Total slides: {len(prs.slides)}")

if __name__ == "__main__":
    main()
